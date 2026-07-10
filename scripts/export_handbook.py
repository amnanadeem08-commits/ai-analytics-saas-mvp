#!/usr/bin/env python3
"""Export Engineering Handbook to PDF, DOCX, and PPTX."""
from __future__ import annotations

from pathlib import Path

ROOT = Path(__file__).resolve().parents[1]
DOC = ROOT / "documentation"
OUT = DOC / "exports"
OUT.mkdir(parents=True, exist_ok=True)


def collect_markdown() -> list[tuple[str, str]]:
    files: list[tuple[str, str]] = []
    preferred = [
        "README.md",
        "01_executive_summary/README.md",
        "02_architecture/README.md",
        "02_architecture/COMPONENT_CATALOG.md",
        "03_folders/README.md",
        "04_workflows/README.md",
        "05_data_science/README.md",
        "06_statistics/README.md",
        "07_power_bi/README.md",
        "08_ai/README.md",
        "09_api/README.md",
        "10_database/README.md",
        "11_frontend/README.md",
        "12_pictorial_evidence/README.md",
        "13_testing/README.md",
        "14_sprints/README.md",
        "15_learning/README.md",
        "16_interview/README.md",
        "17_portfolio/README.md",
        "18_release/README.md",
        "18_release/CONTRIBUTING.md",
        "19_formats/README.md",
        "20_standards/README.md",
    ]
    for rel in preferred:
        p = DOC / rel
        if p.exists():
            files.append((rel, p.read_text(encoding="utf-8")))
    return files


def export_pdf(sections: list[tuple[str, str]]) -> Path:
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    path = OUT / "DataBotAI_v1_Engineering_Handbook.pdf"
    doc = SimpleDocTemplate(str(path), pagesize=LETTER)
    styles = getSampleStyleSheet()
    story = []
    story.append(Paragraph("Data Bot AI v1.0 — Engineering Handbook", styles["Title"]))
    story.append(Spacer(1, 12))
    for rel, text in sections:
        story.append(Paragraph(rel.replace("/", " / "), styles["Heading2"]))
        for line in text.splitlines():
            line = line.strip()
            if not line or line.startswith("```") or line.startswith("|---"):
                continue
            safe = (
                line.replace("&", "&amp;")
                .replace("<", "&lt;")
                .replace(">", "&gt;")
            )
            if line.startswith("#"):
                story.append(Paragraph(safe.lstrip("# ").strip(), styles["Heading3"]))
            else:
                story.append(Paragraph(safe[:2000], styles["BodyText"]))
            story.append(Spacer(1, 4))
    doc.build(story)
    return path


def export_docx(sections: list[tuple[str, str]]) -> Path:
    # Minimal DOCX (OOXML) without python-docx dependency
    import zipfile
    from xml.sax.saxutils import escape

    path = OUT / "DataBotAI_v1_Engineering_Handbook.docx"
    paragraphs = [
        "<w:p><w:r><w:t>Data Bot AI v1.0 — Engineering Handbook</w:t></w:r></w:p>"
    ]
    for rel, text in sections:
        paragraphs.append(f"<w:p><w:r><w:t>{escape(rel)}</w:t></w:r></w:p>")
        for line in text.splitlines():
            if not line.strip():
                continue
            paragraphs.append(
                f"<w:p><w:r><w:t>{escape(line[:4000])}</w:t></w:r></w:p>"
            )
    document_xml = f"""<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<w:document xmlns:w="http://schemas.openxmlformats.org/wordprocessingml/2006/main">
  <w:body>
    {''.join(paragraphs)}
    <w:sectPr/>
  </w:body>
</w:document>"""
    content_types = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Types xmlns="http://schemas.openxmlformats.org/package/2006/content-types">
  <Default Extension="rels" ContentType="application/vnd.openxmlformats-package.relationships+xml"/>
  <Default Extension="xml" ContentType="application/xml"/>
  <Override PartName="/word/document.xml" ContentType="application/vnd.openxmlformats-officedocument.wordprocessingml.document.main+xml"/>
</Types>"""
    rels = """<?xml version="1.0" encoding="UTF-8" standalone="yes"?>
<Relationships xmlns="http://schemas.openxmlformats.org/package/2006/relationships">
  <Relationship Id="rId1" Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/officeDocument" Target="word/document.xml"/>
</Relationships>"""
    with zipfile.ZipFile(path, "w", compression=zipfile.ZIP_DEFLATED) as zf:
        zf.writestr("[Content_Types].xml", content_types)
        zf.writestr("_rels/.rels", rels)
        zf.writestr("word/document.xml", document_xml)
    return path


def export_pptx(sections: list[tuple[str, str]]) -> Path:
    from pptx import Presentation
    from pptx.util import Inches, Pt

    path = OUT / "DataBotAI_v1_Engineering_Handbook.pptx"
    prs = Presentation()
    # Title
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Data Bot AI v1.0"
    slide.placeholders[1].text = "Engineering Handbook · Portfolio · Learning Guide"

    overview = prs.slides.add_slide(prs.slide_layouts[1])
    overview.shapes.title.text = "Handbook Contents"
    overview.placeholders[1].text = "\n".join(
        [
            "01 Executive Summary",
            "02 Architecture",
            "03 Folders",
            "04 Workflows",
            "05 Data Science",
            "06 Statistics",
            "07 Power BI Mapping",
            "08 AI Handbook",
            "09 API",
            "10 Database",
            "11 Frontend",
            "12 Pictorial Evidence",
            "13 Testing (634 passed @ release)",
            "14 Sprints",
            "15 Learning",
            "16 Interview",
            "17 Portfolio",
            "18 Release",
        ]
    )

    for rel, text in sections[:18]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        title = rel.split("/")[0].replace("_", " ")
        slide.shapes.title.text = title[:80]
        body = []
        for line in text.splitlines():
            line = line.strip()
            if line.startswith("#") or line.startswith("|") or line.startswith("```"):
                continue
            if line:
                body.append(line[:180])
            if len(body) >= 8:
                break
        slide.placeholders[1].text = "\n".join(body) or rel

    prs.save(str(path))
    return path


def main() -> None:
    sections = collect_markdown()
    if not sections:
        raise SystemExit("No markdown found — run generate_handbook_part*.py first")
    pdf = export_pdf(sections)
    docx = export_docx(sections)
    pptx = export_pptx(sections)
    print("PDF ", pdf)
    print("DOCX", docx)
    print("PPTX", pptx)


if __name__ == "__main__":
    main()
