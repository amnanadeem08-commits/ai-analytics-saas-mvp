from __future__ import annotations

import io
from datetime import datetime, timezone
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt

from backend.services.export_render_service import chart_to_png_bytes, filter_charts


SLIDE_W = 13.333
SLIDE_H = 7.5
CONTENT_BOTTOM = 6.65


def _rgb(hex_color: str, fallback: str = "#118DFF") -> RGBColor:
    value = (hex_color or fallback).lstrip("#")
    if len(value) != 6:
        value = fallback.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def _text_chunks(text: str, max_chars: int) -> list[str]:
    words = str(text or "").split()
    if not words:
        return [""]
    lines: list[str] = []
    current = ""
    for word in words:
        candidate = f"{current} {word}".strip()
        if len(candidate) > max_chars and current:
            lines.append(current)
            current = word
        else:
            current = candidate
    if current:
        lines.append(current)
    return lines


def _paginate_items(items: list[str], max_lines: int = 15, max_chars: int = 92) -> list[list[str]]:
    pages: list[list[str]] = []
    current: list[str] = []
    current_lines = 0
    for item in items:
        line_count = max(1, len(_text_chunks(item, max_chars)))
        if current and current_lines + line_count > max_lines:
            pages.append(current)
            current = []
            current_lines = 0
        current.append(item)
        current_lines += line_count
    if current:
        pages.append(current)
    return pages or [[]]


def _add_footer(slide, branding: dict[str, Any], primary: RGBColor, page_number: int) -> None:
    footer = slide.shapes.add_textbox(Inches(0.55), Inches(6.95), Inches(12.2), Inches(0.25))
    frame = footer.text_frame
    frame.text = (
        f"{branding.get('company_name', 'AI Analytics')} | "
        f"{datetime.now(timezone.utc).strftime('%Y-%m-%d %H:%M UTC')} | Slide {page_number}"
    )
    frame.paragraphs[0].font.size = Pt(8)
    frame.paragraphs[0].font.color.rgb = RGBColor(100, 116, 139)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.82), Inches(12.2), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = primary
    line.line.fill.background()


def _add_title(slide, title: str, subtitle: str, color: RGBColor) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(12.2), Inches(0.62))
    title_frame = title_box.text_frame
    title_frame.word_wrap = True
    title_frame.text = title[:130]
    title_frame.paragraphs[0].font.size = Pt(24 if len(title) < 75 else 18)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = color

    sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.02), Inches(11.7), Inches(0.38))
    sub_frame = sub_box.text_frame
    sub_frame.word_wrap = True
    sub_frame.text = subtitle[:160]
    sub_frame.paragraphs[0].font.size = Pt(10)
    sub_frame.paragraphs[0].font.color.rgb = RGBColor(90, 98, 112)


def _add_bullets(slide, items: list[str], x: float, y: float, w: float, h: float, font_size: int = 13) -> None:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_bottom = Inches(0.05)
    frame.margin_top = Inches(0.05)
    frame.clear()
    adjusted = max(9, font_size - max(0, len(items) - 8) // 2)
    for index, item in enumerate(items):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = str(item or "")[:420]
        paragraph.font.size = Pt(adjusted)
        paragraph.space_after = Pt(6)


def _new_slide(prs: Presentation, title: str, subtitle: str, branding: dict[str, Any], primary: RGBColor):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, title, subtitle, primary)
    _add_footer(slide, branding, primary, len(prs.slides))
    return slide


def _add_bullet_slides(prs: Presentation, title: str, subtitle: str, items: list[str], branding: dict[str, Any], primary: RGBColor) -> None:
    pages = _paginate_items(items)
    for index, page_items in enumerate(pages):
        slide_title = title if index == 0 else f"{title} Continued"
        slide = _new_slide(prs, slide_title, subtitle, branding, primary)
        _add_bullets(slide, page_items, 0.75, 1.6, 11.8, CONTENT_BOTTOM - 1.6)


def _add_kpi_card(slide, card: dict[str, Any], x: float, y: float, w: float, color: RGBColor) -> None:
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(248, 250, 252)
    shape.line.color.rgb = RGBColor(215, 222, 232)
    frame = shape.text_frame
    frame.word_wrap = True
    frame.margin_left = Inches(0.12)
    frame.margin_right = Inches(0.12)
    frame.text = str(card.get("label", "KPI"))[:34]
    frame.paragraphs[0].font.size = Pt(9)
    frame.paragraphs[0].font.color.rgb = RGBColor(90, 98, 112)
    p = frame.add_paragraph()
    p.text = str(card.get("value", ""))[:22]
    p.font.size = Pt(18)
    p.font.bold = True
    p.font.color.rgb = color
    p2 = frame.add_paragraph()
    p2.text = str(card.get("reason", card.get("description", "")))[:110]
    p2.font.size = Pt(7)


def build_executive_pptx(
    report: dict[str, Any],
    chart_ids: list[str] | None = None,
    package: str = "executive",
) -> bytes:
    if package == "storyboard":
        return build_storyboard_pptx(report, chart_ids=chart_ids, package=package)
    branding = report.get("branding", {})
    primary = _rgb(branding.get("primary_color", "#118DFF"))
    accent = _rgb(branding.get("accent_color", "#E66C37"))

    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    rendered_charts = [(chart, png) for chart in filter_charts(report, chart_ids) if (png := chart_to_png_bytes(chart, width=1060, height=540))]

    cover = _new_slide(
        prs,
        branding.get("report_title", "Executive Decision Intelligence Report"),
        branding.get("company_name", "AI Analytics"),
        branding,
        primary,
    )
    _add_bullets(
        cover,
        [
            f"Package: {package.title()}",
            f"Dataset: {report.get('dataset_id', '')}",
            f"Rows analyzed: {report.get('overview', {}).get('row_count', 0):,}",
            f"Charts prepared: {len(rendered_charts):,}",
        ],
        0.8,
        2.0,
        11.5,
        3.8,
        16,
    )

    executive = report.get("executive_summary", {})
    _add_bullet_slides(
        prs,
        "Executive Summary",
        "Board-ready narrative",
        [
            executive.get("insight", ""),
            f"Reason: {executive.get('reason', '')}",
            f"Action: {executive.get('action', '')}",
        ],
        branding,
        primary,
    )

    quality = report.get("data_quality_score", {})
    if quality:
        _add_bullet_slides(
            prs,
            "Data Quality",
            "Readiness and reliability",
            [
                f"Score: {quality.get('score', '')} ({quality.get('grade', '')})",
                f"Completeness: {quality.get('completeness_pct', '')}%",
                f"Duplicate rate: {quality.get('duplicate_pct', '')}%",
                f"Interpretation: {quality.get('explanation', '')}",
            ],
            branding,
            primary,
        )

    kpi_slide = _new_slide(prs, "Business Health Overview", "Executive KPI cards", branding, primary)
    for index, card in enumerate(report.get("kpi_cards", [])[:8]):
        row = index // 4
        col = index % 4
        _add_kpi_card(kpi_slide, card, 0.55 + col * 3.18, 1.65 + row * 1.45, 2.95, primary)

    framework_items = []
    for block in executive.get("decision_framework", [])[:6]:
        framework_items.extend(
            [
                f"What happened: {block.get('what_happened', '')}",
                f"Why it happened: {block.get('why_it_happened', '')}",
                f"What to do: {block.get('what_to_do', '')}",
                f"Expected impact: {block.get('expected_impact', '')}",
            ]
        )
    _add_bullet_slides(prs, "Root Cause Analysis", "What / Why / Action", framework_items, branding, primary)

    for chart, png in rendered_charts:
        slide = _new_slide(prs, chart.get("title", "Dashboard Visual"), chart.get("chart_type", "Chart").title(), branding, primary)
        slide.shapes.add_picture(io.BytesIO(png), Inches(0.72), Inches(1.55), width=Inches(11.9), height=Inches(4.82))

    recs = [
        item.get("recommendation", item.get("action", ""))
        for item in executive.get("recommendations", []) + executive.get("action_plan", [])
    ]
    _add_bullet_slides(prs, "Recommendations", "Management action plan", recs[:12], branding, primary)

    story = report.get("business_story", {})
    closing = _new_slide(prs, "Closing Summary", "Executive story mode", branding, primary)
    _add_bullets(
        closing,
        [
            f"What happened: {story.get('data_story', '')}",
            f"Expected impact: {story.get('trend_story', '')}",
            f"Executive narrative: {story.get('business_story', '')}",
        ],
        0.75,
        1.65,
        11.8,
        4.8,
    )
    accent_line = closing.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(6.62), Inches(11.8), Inches(0.06))
    accent_line.fill.solid()
    accent_line.fill.fore_color.rgb = accent
    accent_line.line.fill.background()

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()

def _storyboard_section(payload: dict[str, Any], section_id: str) -> dict[str, Any]:
    return next((section for section in payload.get("sections", []) if section.get("section_id") == section_id), {})


def build_storyboard_pptx(
    report: dict[str, Any],
    chart_ids: list[str] | None = None,
    package: str = "storyboard",
) -> bytes:
    branding = report.get("branding", {})
    storyboard = report.get("executive_storyboard", {})
    primary = _rgb(branding.get("primary_color", "#118DFF"))
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)

    cover = _new_slide(prs, "Executive Storyboard", f"Dataset: {report.get('dataset_id', '')} | Package: {package.title()}", branding, primary)
    source = storyboard.get("source_payloads", {})
    _add_bullets(cover, [f"Data Insights: {source.get('data_insights_status', 'unknown')}", f"AI Business Insights: {source.get('ai_business_insights_status', 'unknown')}", f"Dashboard: {source.get('dashboard_status', 'unknown')}"], 0.8, 1.8, 11.5, 3.6, 16)

    summary = _storyboard_section(storyboard, "executive_summary").get("content", {})
    readiness = summary.get("dataset_readiness", {})
    slide = _new_slide(prs, "Executive Summary", "Readiness, business health, opportunity, and risk", branding, primary)
    _add_bullets(slide, [f"Dataset readiness: {readiness.get('score', 0)}/100", f"Overall business health: {summary.get('overall_business_health', 0)}/100", f"Summary: {summary.get('executive_summary', '')}", f"Top opportunity: {summary.get('top_opportunity', '')}", f"Biggest risk: {summary.get('biggest_risk', '')}"], 0.75, 1.55, 11.8, 4.9, 13)

    kpi_slide = _new_slide(prs, "KPI Overview", "Validated KPI cards", branding, primary)
    kpis = _storyboard_section(storyboard, "kpi_overview").get("kpis", [])
    if kpis:
        for index, card in enumerate(kpis[:8]):
            row = index // 4
            col = index % 4
            _add_kpi_card(kpi_slide, card, 0.55 + col * 3.18, 1.65 + row * 1.45, 2.95, primary)
    else:
        _add_bullets(kpi_slide, ["No KPI cards are available for this dataset."], 0.8, 2.0, 11.5, 2.0)

    cards = _storyboard_section(storyboard, "ai_business_insights").get("cards", [])
    insight_items = [f"{card.get('type', 'Insight')}: {card.get('title', '')} | {card.get('executive_recommendation', '')}" for card in cards]
    _add_bullet_slides(prs, "AI Business Insights", "Existing Phase 4 cards", insight_items or ["No AI Business Insight cards are available."], branding, primary)

    charts = _storyboard_section(storyboard, "executive_charts").get("charts", [])
    if chart_ids:
        selected = set(chart_ids)
        charts = [chart for chart in charts if chart.get("chart_id") in selected]
    rendered_charts = [(chart, png) for chart in charts if (png := chart_to_png_bytes(chart, width=1060, height=540))]
    if rendered_charts:
        for chart, png in rendered_charts:
            slide = _new_slide(prs, chart.get("title", "Executive Chart"), chart.get("chart_type", "Chart").title(), branding, primary)
            slide.shapes.add_picture(io.BytesIO(png), Inches(0.72), Inches(1.55), width=Inches(11.9), height=Inches(4.82))
    else:
        slide = _new_slide(prs, "Executive Charts", "Visualization engine output", branding, primary)
        _add_bullets(slide, ["No chart-ready visuals were generated for this dataset."], 0.8, 2.1, 11.5, 2.0)

    recs = _storyboard_section(storyboard, "executive_recommendations").get("recommendations", [])
    rec_items = [f"{rec.get('priority', 'Medium')} priority | Difficulty: {rec.get('difficulty', '')} | Impact: {rec.get('expected_impact', '')} | {rec.get('recommendation', '')}" for rec in recs]
    _add_bullet_slides(prs, "Executive Recommendations", "Ranked action plan", rec_items or ["No ranked recommendations are available."], branding, primary)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue()