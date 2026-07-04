from __future__ import annotations

import html
import io
import json
import re
from datetime import date, datetime
from urllib.parse import urlencode

import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import requests
import streamlit as st

from frontend.api_client.backend_client import BackendClient, DEFAULT_API_BASE_URL
from frontend.components.chart_components import (
    render_plotly_chart_specs,
    render_time_trends,
    render_top_categories,
)
from frontend.components.insight_cards import render_insight
from frontend.components.metric_cards import render_summary_metrics
from frontend.components.ai_insight_panel import render_business_insights_overview
from frontend.components.storyboard_session import add_storyboard_entry
from frontend.utils.backend_utils import (
    LOCAL_MODE_INFO_MESSAGE,
    _build_local_dataset_id,
    _is_local_dataset_id,
    _warn_backend_unavailable,
    ensure_backend_available,
    get_client,
    is_local_dataset_id,
    render_backend_status,
    safe_table,
)
from frontend.utils.session_state import (
    _ensure_default_local_storyboard,
    _sync_storyboard_keys,
    initialize_session_state,
)
from frontend.utils.theme_manager import (
    DEFAULT_BRANDING,
    THEME_PRESETS,
    _ai_dashboard_css,
    _apply_branding_theme,
    _storyboard_theme_snapshot,
    get_active_branding,
    render_branding_editor,
    render_dashboard_settings,
    render_theme_selector,
)
from frontend.app_pages.dashboard_page import _render_kpi_cards
from frontend.app_pages.dataset_page import _local_active_dataframe
from frontend.utils.kpi_helpers import _local_summary, build_data_anomaly_report, build_default_kpis
from frontend.utils.local_helpers import (
    _figure_png_bytes,
    _local_anomaly_rows,
    _local_default_figures,
    _storyboard_chart_figure,
    _storyboard_slide_text_lines,
    select_dataset,
)

def render_reports(client: BackendClient) -> None:
    st.header("Reports")
    dataset_id = select_dataset(client)
    if not dataset_id:
        return

    if _is_local_dataset_id(dataset_id):
        local_df = _local_active_dataframe(dataset_id)
        if local_df is None:
            st.info("Upload a dataset first from Dataset Preview.")
            return
        st.info(LOCAL_MODE_INFO_MESSAGE)
        summary = _local_summary(local_df)
        cols = st.columns(3)
        cols[0].metric("Rows", f"{summary.get('row_count', 0):,}")
        cols[1].metric("Columns", f"{summary.get('column_count', 0):,}")
        cols[2].metric("Duplicates", f"{summary.get('duplicate_rows', 0):,}")
        st.subheader("Export Package")
        with st.container(border=True):
            _local_report_downloads(local_df, dataset_id)
        return

    try:
        report = client.get_report(dataset_id)
    except requests.RequestException as exc:
        _warn_backend_unavailable("Reports")
        return

    branding = report.get("branding", {})
    if branding:
        st.caption(branding.get("company_name", ""))
        st.subheader(branding.get("report_title", "Executive Report"))

    overview = report.get("overview", {})
    cols = st.columns(3)
    cols[0].metric("Rows", f"{overview.get('row_count', 0):,}")
    cols[1].metric("Columns", f"{overview.get('column_count', 0):,}")
    cols[2].metric("Charts", f"{report.get('chart_count', 0):,}")

    guardrails = report.get("analysis_guardrails", {})
    if guardrails:
        with st.expander("Analysis Readiness", expanded=False):
            st.write(guardrails.get("summary", ""))
            readiness_cols = st.columns(4)
            supports = guardrails.get("supports", {})
            readiness_cols[0].metric("KPI", "Yes" if supports.get("kpi_tracking") else "No")
            readiness_cols[1].metric("Trend", "Yes" if supports.get("trend_analysis") else "No")
            readiness_cols[2].metric("Comparison", "Yes" if supports.get("comparison_analysis") else "No")
            readiness_cols[3].metric("Maps", "Yes" if supports.get("geographic_analysis") else "No")
            for item in guardrails.get("invalid_methods", []):
                st.warning(item)

    executive = report.get("executive_summary", {})
    if executive:
        st.subheader("Executive Summary")
        st.write(f"**Insight:** {executive.get('insight', '')}")
        st.write(f"**Reason:** {executive.get('reason', '')}")
        st.write(f"**Action:** {executive.get('action', '')}")

        business_story = report.get("business_story", {})
        if business_story:
            with st.expander("Business Storytelling Engine", expanded=True):
                st.write(f"**Data Story:** {business_story.get('data_story', '')}")
                st.write(f"**Trend Story:** {business_story.get('trend_story', '')}")
                st.write(f"**Business Story:** {business_story.get('business_story', '')}")

        tabs = st.tabs(["Action Framework", "Findings", "Risks", "Opportunities", "Recommendations", "Action Plan"])
        with tabs[0]:
            st.dataframe(safe_table(executive.get("decision_framework", [])), use_container_width=True)
        with tabs[1]:
            st.dataframe(safe_table(executive.get("key_findings", [])), use_container_width=True)
        with tabs[2]:
            risks = executive.get("risks", [])
            if risks:
                st.dataframe(safe_table(risks), use_container_width=True)
            else:
                st.success("No material risks detected from the current evidence.")
        with tabs[3]:
            st.dataframe(safe_table(executive.get("opportunities", [])), use_container_width=True)
        with tabs[4]:
            st.dataframe(safe_table(executive.get("recommendations", [])), use_container_width=True)
        with tabs[5]:
            st.dataframe(safe_table(executive.get("action_plan", [])), use_container_width=True)

    chart_specs = report.get("chart_specs", [])
    chart_labels = {f"{chart.get('title', chart.get('chart_id'))} ({chart.get('chart_type', 'chart')})": chart.get("chart_id") for chart in chart_specs}
    st.subheader("Export Package")
    package = st.selectbox(
        "Report package",
        ["executive", "board", "dashboard", "selected_visuals"],
        format_func=lambda value: {
            "executive": "Executive Report",
            "board": "Board Report",
            "dashboard": "Complete Dashboard",
            "selected_visuals": "Selected Visuals Only",
        }[value],
    )
    selected_labels = st.multiselect(
        "Visuals to include",
        list(chart_labels),
        default=list(chart_labels) if package != "selected_visuals" else list(chart_labels)[: min(3, len(chart_labels))],
    )
    selected_chart_ids = [chart_labels[label] for label in selected_labels if chart_labels.get(label)]
    if chart_specs:
        st.success(
            f"{len(selected_chart_ids)} visual(s) selected. Use the buttons below to download them as PDF, PPTX, PNG, JSON, or CSV."
        )
    else:
        st.info("No generated dashboard visuals are available for this dataset yet.")

    with st.container(border=True):
        render_export_downloads(
            client,
            dataset_id,
            selected_chart_ids,
            package,
            label_prefix="Download",
        )
def render_export_downloads(
    client: BackendClient,
    dataset_id: str,
    chart_ids: list[str] | None = None,
    kpi_ids: list[str] | None = None,
    package: str = "executive",
    label_prefix: str = "",
) -> None:
    if is_local_dataset_id(dataset_id):
        local_df = _local_active_dataframe(dataset_id)
        if local_df is None:
            st.info("Upload a dataset first from Dataset Preview.")
            return
        _local_report_downloads(local_df, dataset_id)
        return

    def export_url(report_format: str) -> str:
        params: list[tuple[str, str]] = [("format", report_format), ("package", package)]
        for chart_id in chart_ids or []:
            params.append(("chart_ids", chart_id))
        for kpi_id in kpi_ids or []:
            params.append(("kpi_ids", kpi_id))
        return f"{client.base_url}/report/{dataset_id}/export?{urlencode(params)}"

    selected_count = len(chart_ids or [])
    kpi_count = len(kpi_ids or [])
    if not chart_ids and not kpi_ids:
        target = "complete dashboard"
    else:
        target = f"{selected_count} visual{'s' if selected_count != 1 else ''} and {kpi_count} KPI{'s' if kpi_count != 1 else ''}"
    st.markdown(f"**Download exports for {target}**")
    st.caption(
        "Click one format to generate and download it. Large PDF/PPTX exports may take a moment, but the page will not time out."
    )
    download_path = "C:\\Users\\DELL\\Downloads"
    file_names = {
        "JSON": f"{dataset_id}_report.json",
        "CSV": f"{dataset_id}.csv",
        "PDF": f"{dataset_id}_executive_report.pdf",
        "PPTX": f"{dataset_id}_executive_deck.pptx",
        "Excel": f"{dataset_id}_executive_report.xlsx",
        "PNG": f"{dataset_id}_dashboard_snapshot.png",
    }
    st.info(
        f"Download location: browser Downloads folder, usually `{download_path}` on this machine. "
        "If your browser asks where to save files, it will use the folder you choose."
    )
    with st.expander("Export file names"):
        for label, file_name in file_names.items():
            st.write(f"**{label}:** `{file_name}`")

    safe_prefix = f"{label_prefix} " if label_prefix else ""
    col1, col2, col3, col4, col5, col6 = st.columns(6)
    col1.link_button(
        f"{safe_prefix}JSON",
        export_url("json"),
        use_container_width=True,
    )
    col2.link_button(
        f"{safe_prefix}CSV",
        export_url("csv"),
        use_container_width=True,
    )
    col3.link_button(
        f"{safe_prefix}PDF",
        export_url("pdf"),
        use_container_width=True,
    )
    col4.link_button(
        f"{safe_prefix}PPTX",
        export_url("pptx"),
        use_container_width=True,
    )
    col5.link_button(
        f"{safe_prefix}Excel",
        export_url("xlsx"),
        use_container_width=True,
    )
    col6.link_button(
        f"{safe_prefix}PNG",
        export_url("png"),
        use_container_width=True,
    )
def render_presentation_mode(client: BackendClient) -> None:
    st.header("Presentation Mode")
    dataset_id = select_dataset(client)
    if not dataset_id:
        return

    if is_local_dataset_id(dataset_id):
        st.info("Presentation Mode requires backend connection. Local dataset dashboards and previews are available.")
        return

    try:
        report = client.get_report(dataset_id)
    except requests.RequestException as exc:
        _warn_backend_unavailable("Presentation Mode")
        return

    slides = _presentation_slides(report)
    index = st.slider("Slide", 1, len(slides), 1) - 1
    slide = slides[index]
    branding = report.get("branding", {})
    theme = report.get("theme", {})
    primary = branding.get("primary_color", theme.get("primary", "var(--brand-primary)"))

    st.markdown(
        f"""
        <style>
        .presentation-frame {{
            min-height: 620px;
            border-radius: 10px;
            padding: 38px 44px;
            background: {theme.get('surface', 'var(--ui-surface)')};
            border: 1px solid {theme.get('border', 'var(--surface-border)')};
            box-shadow: 0 18px 42px rgba(15, 23, 42, 0.10);
        }}
        .presentation-title {{
            color: {primary};
            font-size: 2rem;
            font-weight: 800;
            margin-bottom: 8px;
        }}
        .presentation-subtitle {{
            color: {theme.get('muted_text', 'var(--text-muted)')};
            margin-bottom: 28px;
        }}
        .presentation-body {{
            color: {theme.get('text', 'var(--text-color)')};
            font-size: 1.05rem;
            line-height: 1.55;
            margin-bottom: 14px;
        }}
        </style>
        """,
        unsafe_allow_html=True,
    )
    with st.container(border=True):
        st.markdown(f'<div class="presentation-title">{slide["title"]}</div>', unsafe_allow_html=True)
        st.markdown(f'<div class="presentation-subtitle">{branding.get("company_name", "AI Analytics")} | Slide {index + 1} of {len(slides)}</div>', unsafe_allow_html=True)
        if slide.get("kpis"):
            _render_kpi_cards(slide["kpis"], theme, key_prefix="executive_dashboard")
        for item in slide.get("body", []):
            if item:
                st.markdown(f'<div class="presentation-body">{item}</div>', unsafe_allow_html=True)

    st.subheader("Download Presentation")
    chart_specs = report.get("chart_specs", [])
    chart_labels = {
        f"{chart.get('title', chart.get('chart_id'))} ({chart.get('chart_type', 'chart')})": chart.get("chart_id")
        for chart in chart_specs
    }
    selected_labels = st.multiselect(
        "Visuals to include in presentation exports",
        list(chart_labels),
        default=list(chart_labels),
        key="presentation_export_visuals",
    )
    selected_chart_ids = [chart_labels[label] for label in selected_labels if chart_labels.get(label)]
    with st.container(border=True):
        render_export_downloads(
            client,
            dataset_id,
            selected_chart_ids,
            "dashboard",
            label_prefix="Download",
        )
def _presentation_slides(report: dict) -> list[dict]:
    executive = report.get("executive_summary", {})
    story = report.get("business_story", {})
    kpis = report.get("kpi_cards", [])
    blocks = executive.get("decision_framework", [])
    return [
        {"title": "Executive Summary", "body": [executive.get("insight", ""), executive.get("reason", ""), executive.get("action", "")]},
        {"title": "Business Health Overview", "body": [story.get("business_story", "")]},
        {"title": "Key KPIs", "kpis": kpis[:4], "body": []},
        {"title": "Revenue Analysis", "body": [blocks[0].get("what_happened", "") if blocks else "", blocks[0].get("why_it_happened", "") if blocks else ""]},
        {"title": "Customer Analysis", "body": [blocks[1].get("what_happened", "") if len(blocks) > 1 else "Segment analysis appears when customer or segment fields are available."]},
        {"title": "Root Cause Analysis", "body": [block.get("why_it_happened", "") for block in blocks[:3]]},
        {"title": "Risks", "body": [item.get("risk", "") + ": " + item.get("why_it_matters", "") for item in executive.get("risks", [])] or ["No material risks detected from current evidence."]},
        {"title": "Opportunities", "body": [item.get("opportunity", "") + ": " + item.get("why", "") for item in executive.get("opportunities", [])]},
        {"title": "Recommendations", "body": [item.get("recommendation", "") for item in executive.get("recommendations", [])]},
        {"title": "Action Plan", "body": [item.get("action", "") for item in executive.get("action_plan", [])]},
    ]
def _local_report_downloads(df: pd.DataFrame, dataset_id: str, storyboard_items: list[dict] | None = None) -> None:
    with st.spinner("Preparing local export package with real dashboard visuals..."):
        files = _local_export_files(df, dataset_id, storyboard_items=storyboard_items)

    col1, col2, col3, col4, col5, col6 = st.columns(6)
    col1.download_button("Download JSON", data=files["json"], file_name=f"{dataset_id}_report.json", mime="application/json", use_container_width=True)
    col2.download_button("Download CSV", data=files["csv"], file_name=f"{dataset_id}.csv", mime="text/csv", use_container_width=True)
    col3.download_button("Download PDF", data=files["pdf"], file_name=f"{dataset_id}_executive_report.pdf", mime="application/pdf", use_container_width=True)
    col4.download_button("Download PPTX", data=files["pptx"], file_name=f"{dataset_id}_executive_deck.pptx", mime="application/vnd.openxmlformats-officedocument.presentationml.presentation", use_container_width=True)
    col5.download_button("Download Excel", data=files["xlsx"], file_name=f"{dataset_id}_executive_report.xlsx", mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet", use_container_width=True)
    col6.download_button("Download PNG", data=files["png"], file_name=f"{dataset_id}_dashboard_snapshot.png", mime="image/png", use_container_width=True, disabled=not bool(files["png"]))
    st.caption("Local exports include automatically generated visuals when no manual storyboard/report visuals are selected.")
def _local_export_files(df: pd.DataFrame, dataset_id: str, storyboard_items: list[dict] | None = None) -> dict[str, bytes]:
    base_figures = _local_default_figures(df)
    if storyboard_items:
        figures: list[dict] = []
        seen_chart_keys: set[str] = set()
        for slide in storyboard_items:
            for chart_ref in slide.get("charts", []):
                chart_key = str(chart_ref.get("chart_id") or chart_ref.get("title") or "")
                if chart_key and chart_key in seen_chart_keys:
                    continue
                figure, title, kind = _storyboard_chart_figure(df, chart_ref, base_figures)
                if figure is None:
                    continue
                figures.append({"title": title, "figure": figure, "kind": kind})
                if chart_key:
                    seen_chart_keys.add(chart_key)
        if not figures:
            figures = base_figures
    else:
        figures = base_figures

    summary = _local_summary(df)
    kpis = build_default_kpis(df, dataset_id)
    anomaly_report = build_data_anomaly_report(df)
    files: dict[str, bytes] = {
        "json": json.dumps(
            {
                "summary": summary,
                "kpi_cards": kpis,
                "anomaly_report": anomaly_report,
                "visuals": [item["title"] for item in figures],
                "storyboard": [
                    {
                        "slide_id": item.get("slide_id"),
                        "title": item.get("title"),
                        "section_type": item.get("section_type"),
                    }
                    for item in (storyboard_items or [])
                ],
            },
            indent=2,
        ).encode("utf-8"),
        "csv": df.to_csv(index=False).encode("utf-8"),
    }

    excel_buffer = io.BytesIO()
    with pd.ExcelWriter(excel_buffer, engine="openpyxl") as writer:
        df.head(5000).to_excel(writer, sheet_name="Raw Data", index=False)
        pd.DataFrame([summary]).to_excel(writer, sheet_name="Summary", index=False)
        pd.DataFrame(kpis).to_excel(writer, sheet_name="Executive Summary", index=False)
        pd.DataFrame([{"Column": col, "Type": str(dtype), "Missing": int(df[col].isna().sum()), "Unique": int(df[col].nunique(dropna=True))} for col, dtype in df.dtypes.items()]).to_excel(writer, sheet_name="Column Schema", index=False)
        pd.DataFrame(anomaly_report or [{"message": "No major statistical anomalies detected in this dataset."}]).to_excel(writer, sheet_name="Cleaning Anomaly Report", index=False)
        pd.DataFrame([{"Visual": item["title"], "Type": item["kind"]} for item in figures]).to_excel(writer, sheet_name="Dashboard Visuals", index=False)
        workbook = writer.book
        visual_sheet = workbook.create_sheet("Visuals")
        row = 1
        from openpyxl.drawing.image import Image as ExcelImage
        for index, item in enumerate(figures[:6], start=1):
            png = _figure_png_bytes(item["figure"])
            image_buffer = io.BytesIO(png)
            image = ExcelImage(image_buffer)
            image.width = 640
            image.height = 380
            visual_sheet.cell(row=row, column=1, value=item["title"])
            visual_sheet.add_image(image, f"A{row + 1}")
            row += 24
    files["xlsx"] = excel_buffer.getvalue()

    from reportlab.lib.pagesizes import landscape, letter
    from reportlab.lib.utils import ImageReader
    from reportlab.pdfgen import canvas
    pdf_buffer = io.BytesIO()
    pdf = canvas.Canvas(pdf_buffer, pagesize=landscape(letter))
    width, height = landscape(letter)
    pdf.setTitle(f"{dataset_id} executive report")
    pdf.setFont("Helvetica-Bold", 18)
    pdf.drawString(40, height - 42, "Executive Analytics Export")
    pdf.setFont("Helvetica", 10)
    pdf.drawString(40, height - 60, f"Rows: {summary.get('row_count', 0):,}   Columns: {summary.get('column_count', 0):,}   Visuals: {len(figures):,}")
    y = height - 92
    pdf.setFont("Helvetica-Bold", 13)
    pdf.drawString(48, y, "KPI Overview")
    y -= 18
    pdf.setFont("Helvetica", 9)
    for card in kpis[:6]:
        if y < 120:
            break
        pdf.drawString(54, y, f"- {card.get('title') or card.get('label')}: {card.get('formatted_value') or card.get('value')} | {card.get('short_interpretation', '')}"[:145])
        y -= 13
    y -= 8
    if anomaly_report:
        pdf.setFont("Helvetica-Bold", 13)
        pdf.drawString(48, y, "Anomalies and Risks")
        y -= 16
        pdf.setFont("Helvetica", 9)
        for item in anomaly_report[:4]:
            if y < 80:
                break
            pdf.drawString(54, y, f"- {item['severity'].title()} {item['issue_type']} in {item['column']}: {item['explanation']}"[:145])
            y -= 13
    if storyboard_items:
        for line in _storyboard_slide_text_lines(storyboard_items[0])[:6]:
            pdf.drawString(48, y, line[:130])
            y -= 14
    else:
        for item in _local_anomaly_rows(df)[:5]:
            pdf.drawString(48, y, f"- {item[:120]}")
            y -= 14
    pdf.showPage()

    if storyboard_items:
        for slide in storyboard_items:
            pdf.setFont("Helvetica-Bold", 15)
            pdf.drawString(40, height - 38, str(slide.get("title") or "Storyboard Slide")[:90])
            pdf.setFont("Helvetica", 10)
            y = height - 62
            for line in _storyboard_slide_text_lines(slide):
                if y < 220:
                    break
                pdf.drawString(46, y, line[:130])
                y -= 14
            chart_ref = next((item for item in slide.get("charts", []) if isinstance(item, dict)), None)
            if chart_ref is not None:
                figure, _, _ = _storyboard_chart_figure(df, chart_ref, base_figures)
                if figure is not None:
                    png = _figure_png_bytes(figure)
                    pdf.drawImage(ImageReader(io.BytesIO(png)), 42, 42, width=width - 84, height=170, preserveAspectRatio=True, anchor="c")
            pdf.showPage()
    else:
        for item in figures[:8]:
            png = _figure_png_bytes(item["figure"])
            pdf.setFont("Helvetica-Bold", 15)
            pdf.drawString(40, height - 38, item["title"][:90])
            pdf.drawImage(ImageReader(io.BytesIO(png)), 42, 54, width=width - 84, height=height - 112, preserveAspectRatio=True, anchor="c")
            pdf.showPage()
    pdf.save()
    files["pdf"] = pdf_buffer.getvalue()

    from pptx import Presentation
    from pptx.util import Inches, Pt
    deck = Presentation()
    title_slide = deck.slides.add_slide(deck.slide_layouts[5])
    title_slide.shapes.title.text = "Executive Analytics Export"
    textbox = title_slide.shapes.add_textbox(Inches(0.7), Inches(1.6), Inches(11.8), Inches(1.2))
    text_frame = textbox.text_frame
    text_frame.text = f"{summary.get('row_count', 0):,} rows, {summary.get('column_count', 0):,} columns, {len(figures):,} generated visuals"
    text_frame.paragraphs[0].font.size = Pt(18)
    kpi_slide = deck.slides.add_slide(deck.slide_layouts[5])
    kpi_slide.shapes.title.text = "Business Health Overview"
    for idx, card in enumerate(kpis[:8]):
        left = Inches(0.7 + (idx % 4) * 3.1)
        top = Inches(1.25 + (idx // 4) * 2.15)
        box = kpi_slide.shapes.add_textbox(left, top, Inches(2.85), Inches(1.65))
        tf = box.text_frame
        tf.clear()
        tf.text = str(card.get("title") or card.get("label") or "KPI")[:34]
        tf.paragraphs[0].font.size = Pt(12)
        tf.paragraphs[0].font.bold = True
        p = tf.add_paragraph()
        p.text = str(card.get("formatted_value") or card.get("value") or "")[:32]
        p.font.size = Pt(22)
        p.font.bold = True
        p2 = tf.add_paragraph()
        p2.text = str(card.get("short_interpretation") or card.get("description") or "")[:90]
        p2.font.size = Pt(9)

    if storyboard_items:
        for item in storyboard_items:
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = str(item.get("title") or "Storyboard Slide")[:80]
            body = slide.shapes.add_textbox(Inches(0.7), Inches(1.15), Inches(5.0), Inches(5.9)).text_frame
            body.clear()
            text_lines = _storyboard_slide_text_lines(item)
            if text_lines:
                body.text = text_lines[0]
                for line in text_lines[1:7]:
                    paragraph = body.add_paragraph()
                    paragraph.text = line
            chart_ref = next((chart for chart in item.get("charts", []) if isinstance(chart, dict)), None)
            if chart_ref is not None:
                figure, _, _ = _storyboard_chart_figure(df, chart_ref, base_figures)
                if figure is not None:
                    png = _figure_png_bytes(figure)
                    slide.shapes.add_picture(io.BytesIO(png), Inches(5.9), Inches(1.1), width=Inches(6.7), height=Inches(5.9))
    else:
        for item in figures[:8]:
            slide = deck.slides.add_slide(deck.slide_layouts[5])
            slide.shapes.title.text = item["title"][:80]
            png = _figure_png_bytes(item["figure"])
            slide.shapes.add_picture(io.BytesIO(png), Inches(0.7), Inches(1.15), width=Inches(12.0), height=Inches(5.8))
    ppt_buffer = io.BytesIO()
    deck.save(ppt_buffer)
    files["pptx"] = ppt_buffer.getvalue()
    files["png"] = _figure_png_bytes(figures[0]["figure"]) if figures else b""
    return files
