from __future__ import annotations

import html
import io
import json
import re
from datetime import date, datetime
from urllib.parse import urlencode

import pandas as pd
import plotly.graph_objects as go
import plotly.io as pio
import requests
import streamlit as st

from frontend.api_client.backend_client import BackendClient, DEFAULT_API_BASE_URL
from frontend.components.chart_components import (
    render_plotly_chart_specs,
    render_time_trends,
    render_top_categories,
)
from frontend.components.insight_cards import render_insight
from frontend.components.metric_cards import render_summary_metrics
from frontend.components.ai_insight_panel import render_business_insights_overview
from frontend.components.storyboard_session import add_storyboard_entry
from frontend.utils.backend_utils import (
    _build_local_dataset_id,
    _is_local_dataset_id,
    _warn_backend_unavailable,
    ensure_backend_available,
    get_client,
    is_local_dataset_id,
    render_backend_status,
    safe_table,
)
from frontend.utils.session_state import (
    _ensure_default_local_storyboard,
    _sync_storyboard_keys,
    initialize_session_state,
)
from frontend.utils.theme_manager import (
    DEFAULT_BRANDING,
    THEME_PRESETS,
    _ai_dashboard_css,
    _apply_branding_theme,
    _storyboard_theme_snapshot,
    get_active_branding,
    render_branding_editor,
    render_dashboard_settings,
    render_theme_selector,
)
from frontend.pages.dataset_page import _local_active_dataframe
from frontend.pages.reports_page import render_export_downloads
from frontend.utils.kpi_helpers import build_storyboard_kpis, local_column_groups, _local_kpi_cards, local_summary, quality_score
from frontend.utils.local_helpers import (
    _add_recommended_visual_slide,
    _build_local_visual_schema,
    _dataset_display_name,
    _figure_png_bytes,
    _local_anomaly_rows,
    _local_default_figures,
    _render_local_visual,
    _zscore_outlier_notes,
    select_dataset,
)


def _render_storyboard_dark_theme_css() -> None:
    st.markdown(
        """
        <style>
        .stButton button {
            background: #1D4ED8 !important;
            color: #FFFFFF !important;
            border-radius: 20px !important;
            border: 1px solid #2563EB !important;
            font-weight: 700 !important;
        }
        .story-slide {
            background: #0D1117;
            border-radius: 12px;
            padding: 24px;
            border: 1px solid #1F2937;
            margin: 14px 0 20px;
            color: #E5E7EB;
        }
        .story-header {
            background: #0A0F1E;
            padding: 12px 20px;
            border-bottom: 1px solid #1F2937;
            border-radius: 10px 10px 0 0;
            color: #F8FAFC;
            font-weight: 800;
            letter-spacing: .02em;
            margin: -8px -8px 18px;
        }
        .story-grid {
            display: grid;
            grid-template-columns: minmax(160px, 0.32fr) minmax(0, 0.68fr);
            gap: 18px;
            align-items: stretch;
        }
        .story-kpi-card {
            background: #161B22;
            border-left: 3px solid #0078D4;
            padding: 16px;
            margin: 8px 0;
            border-radius: 8px;
        }
        .story-kpi-value {
            color: #F8FAFC;
            font-size: 1.45rem;
            font-weight: 900;
        }
        .story-kpi-label {
            color: #94A3B8;
            font-size: .82rem;
            margin-top: 4px;
        }
        .story-insight-box {
            background: #0F172A;
            border: 1px solid #334155;
            border-radius: 8px;
            padding: 16px;
            color: #CBD5E1;
            margin-top: 12px;
        }
        .story-recommendation-box {
            background: #052e16;
            border-left: 4px solid #22C55E;
            padding: 16px;
            border-radius: 8px;
            color: #DCFCE7;
            margin-top: 14px;
        }
        .story-section-title {
            color: #93C5FD;
            font-size: .76rem;
            font-weight: 900;
            letter-spacing: .08em;
            text-transform: uppercase;
            margin-bottom: 8px;
        }
        .story-counter {
            color: #64748B;
            text-align: center;
            font-weight: 700;
            padding-top: 8px;
        }
        @media (max-width: 900px) {
            .story-grid { grid-template-columns: 1fr; }
        }
        </style>
        """,
        unsafe_allow_html=True,
    )


def _storyboard_kpi_card_html(label: object, value: object) -> str:
    return (
        '<div class="story-kpi-card">'
        f'<div class="story-kpi-value">{html.escape(str(value))}</div>'
        f'<div class="story-kpi-label">{html.escape(str(label))}</div>'
        '</div>'
    )

def render_storyboard_builder(client: BackendClient) -> None:
    _render_storyboard_dark_theme_css()
    st.header("Storyboard Builder")
    st.caption("Turn Dashboard Studio visuals into a Tableau-style business story for executive review.")
    dataset_id = select_dataset(client)
    if not dataset_id:
        return

    storyboard_key = f"dashboard_studio_storyboard_{dataset_id}"
    slide_key = f"storyboard_slide_{dataset_id}"
    visuals = st.session_state.get(storyboard_key, [])
    local_storyboard_df = _local_active_dataframe(dataset_id) if is_local_dataset_id(dataset_id) else None

    if local_storyboard_df is not None:
        _render_local_storyboard_builder(dataset_id, local_storyboard_df)
        return

    if not visuals:
        st.info("No storyboard visuals were found for this backend dataset. Add visuals from Dashboard Studio to begin this storyboard.")
        if local_storyboard_df is not None and st.button("Auto-build storyboard from dashboard visuals", type="primary", use_container_width=True):
            st.session_state[storyboard_key] = _auto_storyboard_entries(local_storyboard_df, dataset_id)
            st.rerun()
        elif local_storyboard_df is None:
            st.caption("Backend dashboard visuals can be added from Dashboard Studio when the backend is available.")
        return
    if local_storyboard_df is not None:
        schema = _build_local_visual_schema(local_storyboard_df)
    else:
        try:
            schema = client.get_visual_builder_schema(dataset_id)
        except requests.RequestException:
            schema = {}

    template = st.selectbox(
        "Storyboard template",
        [
            "Executive Overview",
            "Sales Performance Story",
            "Customer Churn Story",
            "Inventory Health Story",
            "Financial Performance Story",
            "Marketing Performance Story",
            "General Business Review",
        ],
    )
    layout_mode = st.selectbox(
        "Slide layout",
        ["Visual + Summary", "Visual only", "Table only", "Summary only", "KPI + Chart", "Full Storyboard"],
    )
    include_options = st.multiselect(
        "Include sections",
        [
            "Executive Summary",
            "KPI Overview",
            "Trend Analysis",
            "Category Comparison",
            "Detailed Table",
            "Recommendations",
            "Risk Analysis",
            "Opportunity Analysis",
            "Location Insights",
        ],
        default=["Executive Summary", "KPI Overview", "Recommendations"],
    )

    slide_count = len(visuals)
    st.session_state.setdefault(slide_key, 0)
    st.session_state[slide_key] = min(st.session_state[slide_key], slide_count - 1)

    st.progress((st.session_state[slide_key] + 1) / slide_count)
    nav_left, nav_mid, nav_right = st.columns([1, 2, 1])
    if nav_left.button("Previous slide", use_container_width=True, disabled=st.session_state[slide_key] == 0):
        st.session_state[slide_key] -= 1
        st.rerun()
    selected_slide = nav_mid.selectbox(
        "Slide selector",
        list(range(1, slide_count + 1)),
        index=st.session_state[slide_key],
    )
    st.session_state[slide_key] = selected_slide - 1
    if nav_right.button("Next slide", use_container_width=True, disabled=st.session_state[slide_key] >= slide_count - 1):
        st.session_state[slide_key] += 1
        st.rerun()

    current = visuals[st.session_state[slide_key]]
    st.markdown(f'<div class="story-counter">{st.session_state[slide_key] + 1} of {slide_count}</div>', unsafe_allow_html=True)

    with st.container():
        st.markdown(f'<div class="story-slide"><div class="story-header">:material/slideshow: Slide {st.session_state[slide_key] + 1}/{slide_count} &nbsp; | &nbsp; {html.escape(str(current.get("title", "Storyboard Slide")))}</div><div class="story-grid">', unsafe_allow_html=True)
        left_panel, right_panel = st.columns([0.32, 0.68])
        left_panel.markdown('<div class="story-section-title">KPI Cards</div>', unsafe_allow_html=True)

        storyboard_kpis = build_storyboard_kpis(schema)
        if layout_mode in {"KPI + Chart", "Full Storyboard"}:
            for kpi in storyboard_kpis:
                left_panel.markdown(_storyboard_kpi_card_html(kpi["label"], kpi["value"]), unsafe_allow_html=True)

        insight = current.get("short_ai_insight") or current.get("business_meaning", "")
        explanation = current.get("business_meaning") or "This slide summarizes the selected dashboard visual for business review."
        recommendation = current.get("why_useful") or "Use this slide to guide the next management discussion."

        spec = current.get("spec", {})
        if spec.get("chart_type") == "kpi":
            st.metric(spec.get("title", "KPI"), spec.get("value", "—"))
        elif layout_mode not in {"Summary only", "Table only"} and spec.get("dimension"):
            if local_storyboard_df is not None:
                visual = _render_local_visual(local_storyboard_df, spec)
                chart = visual.get("chart", {})
                plotly_spec = chart.get("plotly", {})
                with right_panel:
                    with right_panel:
                        st.plotly_chart(go.Figure(data=plotly_spec.get("data", []), layout=plotly_spec.get("layout", {})), use_container_width=True)
            else:
                try:
                    visual = client.render_visual(dataset_id, spec)
                    chart = visual.get("chart", {})
                    plotly_spec = chart.get("plotly", {})
                    st.plotly_chart(go.Figure(data=plotly_spec.get("data", []), layout=plotly_spec.get("layout", {})), use_container_width=True)
                except requests.RequestException:
                    _warn_backend_unavailable("Storyboard visual rendering")
        elif layout_mode == "Table only":
            with right_panel:
                st.dataframe(pd.DataFrame([current.get("spec", {})]), use_container_width=True)

        if layout_mode in {"Visual + Summary", "Summary only", "KPI + Chart", "Full Storyboard"}:
            right_panel.markdown(
                f'<div class="story-insight-box"><div class="story-section-title">Insight</div>{html.escape(str(insight or explanation or "Insight will appear when more evidence is available."))}</div>',
                unsafe_allow_html=True,
            )
            right_panel.markdown(
                f'<div class="story-recommendation-box"><b>💡 Recommendation</b><br>{html.escape(str(recommendation))}</div>',
                unsafe_allow_html=True,
            )
        st.markdown('</div></div>', unsafe_allow_html=True)

    storyboard_chart_ids = [item.get("chart_id") for item in visuals if item.get("chart_id") and not str(item.get("chart_id")).startswith("kpi_")]
    storyboard_kpi_ids = [item.get("kpi_id") for item in visuals if item.get("kpi_id")]
    with st.container(border=True):
        st.subheader("Storyboard Export")
        st.caption("Export only visuals currently in this storyboard while preserving active branding.")
        render_export_downloads(
            client,
            dataset_id,
            chart_ids=storyboard_chart_ids,
            kpi_ids=storyboard_kpi_ids,
            package="board",
            label_prefix="Storyboard",
        )
def _render_local_storyboard_builder(dataset_id: str, df: pd.DataFrame) -> None:
    items = _ensure_default_local_storyboard(dataset_id, df)

    st.success("Auto-built executive storyboard generated from your selected dataset.")
    action_cols = st.columns(5)
    if action_cols[0].button("Regenerate Storyboard", use_container_width=True):
        built = build_default_storyboard(df, dataset_id, st.session_state.get("selected_theme"), dict(st.session_state))
        _sync_storyboard_keys(dataset_id, built)
        st.session_state["storyboard_user_edited"] = False
        st.rerun()
    if action_cols[1].button("Clear Storyboard", use_container_width=True):
        _sync_storyboard_keys(dataset_id, [])
        st.session_state["storyboard_user_edited"] = True
        st.rerun()
    if action_cols[2].button("Add Recommended Visuals", use_container_width=True):
        updated = _add_recommended_visual_slide(items, df, dataset_id)
        _sync_storyboard_keys(dataset_id, updated)
        st.session_state["storyboard_user_edited"] = True
        st.rerun()

    with action_cols[3]:
        st.caption("Export PPTX")
    with action_cols[4]:
        st.caption("Export PDF")

    if not items:
        st.info("Storyboard is currently empty. Use Regenerate Storyboard to create the default executive board.")
        return

    st.progress(1 / max(len(items), 1))

    remove_index: int | None = None
    move_up_index: int | None = None
    move_down_index: int | None = None
    updated_items = [dict(item) for item in items]

    for idx, slide in enumerate(updated_items):
        slide_id = slide.get("slide_id", f"slide_{idx+1}")
        with st.container():
            st.markdown(f'<div class="story-slide"><div class="story-header">:material/slideshow: Slide {idx + 1}/{len(updated_items)} &nbsp; | &nbsp; {html.escape(str(slide.get("title", "Storyboard Slide")))}</div><div class="story-grid">', unsafe_allow_html=True)
            header_cols = st.columns([3.2, 1.2, 1.2, 1.2])
            new_title = header_cols[0].text_input(
                "Slide title",
                value=slide.get("title", "Storyboard Slide"),
                key=f"story_title_{slide_id}",
                label_visibility="collapsed",
            )
            if new_title != slide.get("title", "Storyboard Slide"):
                slide["title"] = new_title
                st.session_state["storyboard_user_edited"] = True

            if header_cols[1].button("Up", key=f"story_up_{slide_id}", use_container_width=True, disabled=idx == 0):
                move_up_index = idx
            if header_cols[2].button("Down", key=f"story_down_{slide_id}", use_container_width=True, disabled=idx >= len(updated_items) - 1):
                move_down_index = idx
            if header_cols[3].button("Remove", key=f"story_remove_{slide_id}", use_container_width=True):
                remove_index = idx

            st.markdown(f'<div class="story-counter">{idx + 1} of {len(updated_items)}</div>', unsafe_allow_html=True)
            st.caption(slide.get("content", {}).get("description", ""))

            kpi_names = [kpi.get("label", "KPI") for kpi in slide.get("kpis", [])]
            chart_names = [chart.get("title", "Chart") for chart in slide.get("charts", [])]
            insight_names = slide.get("insights", [])
            left_panel, right_panel = st.columns([0.32, 0.68])
            left_panel.markdown('<div class="story-section-title">KPI Cards</div>', unsafe_allow_html=True)
            if slide.get("kpis"):
                for kpi in slide["kpis"]:
                    left_panel.markdown(_storyboard_kpi_card_html(kpi.get("label", "KPI"), kpi.get("value", "-")), unsafe_allow_html=True)
            else:
                left_panel.markdown('<div class="story-kpi-card"><div class="story-kpi-label">No KPIs on this slide</div></div>', unsafe_allow_html=True)

            first_chart = next((chart for chart in slide.get("charts", []) if chart.get("plotly", {}).get("data")), None)
            if first_chart:
                try:
                    fig = go.Figure(
                        data=first_chart.get("plotly", {}).get("data", []),
                        layout=first_chart.get("plotly", {}).get("layout", {}),
                    )
                    with right_panel:
                        st.plotly_chart(fig, use_container_width=True)
                except Exception:
                    right_panel.caption("Chart rendering skipped because this chart payload is not valid in the current session.")
            else:
                right_panel.markdown(
                    f'<div class="story-insight-box"><div class="story-section-title">Charts</div>{html.escape(", ".join(chart_names[:6]) if chart_names else "No chart on this slide")}</div>',
                    unsafe_allow_html=True,
                )
            right_panel.markdown(
                f'<div class="story-insight-box"><div class="story-section-title">Insight</div>{html.escape(", ".join([str(item) for item in insight_names[:3]]) if insight_names else slide.get("content", {}).get("description", "Insight will appear when available."))}</div>',
                unsafe_allow_html=True,
            )
            right_panel.markdown(
                '<div class="story-recommendation-box"><b>💡 Recommendation</b><br>Use this slide to guide the next executive discussion.</div>',
                unsafe_allow_html=True,
            )
            st.markdown('</div></div>', unsafe_allow_html=True)

    if remove_index is not None:
        updated_items.pop(remove_index)
        _sync_storyboard_keys(dataset_id, updated_items)
        st.session_state["storyboard_user_edited"] = True
        st.rerun()
    if move_up_index is not None and move_up_index > 0:
        updated_items[move_up_index - 1], updated_items[move_up_index] = updated_items[move_up_index], updated_items[move_up_index - 1]
        _sync_storyboard_keys(dataset_id, updated_items)
        st.session_state["storyboard_user_edited"] = True
        st.rerun()
    if move_down_index is not None and move_down_index < len(updated_items) - 1:
        updated_items[move_down_index + 1], updated_items[move_down_index] = updated_items[move_down_index], updated_items[move_down_index + 1]
        _sync_storyboard_keys(dataset_id, updated_items)
        st.session_state["storyboard_user_edited"] = True
        st.rerun()

    _sync_storyboard_keys(dataset_id, updated_items)

    with st.container(border=True):
        st.subheader("Storyboard Export")
        st.caption("Exports use the current storyboard items and do not require manual visual selection.")
        _render_local_storyboard_exports(updated_items, dataset_id)
def _render_local_storyboard_exports(items: list[dict], dataset_id: str) -> None:
    export_files = _local_storyboard_export_files(items, dataset_id)
    action_cols = st.columns(3)
    action_cols[0].download_button(
        "Export PPTX",
        data=export_files.get("pptx", b""),
        file_name=f"{dataset_id}_storyboard.pptx",
        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        use_container_width=True,
        disabled=not bool(export_files.get("pptx")),
    )
    action_cols[1].download_button(
        "Export PDF",
        data=export_files.get("pdf", b""),
        file_name=f"{dataset_id}_storyboard.pdf",
        mime="application/pdf",
        use_container_width=True,
        disabled=not bool(export_files.get("pdf")),
    )
    action_cols[2].download_button(
        "Export JSON",
        data=export_files.get("json", b"{}"),
        file_name=f"{dataset_id}_storyboard.json",
        mime="application/json",
        use_container_width=True,
    )
def build_default_storyboard(df: pd.DataFrame, dataset_id: str, theme: str | None, settings: dict | None) -> list[dict]:
    summary = local_summary(df)
    numeric, categorical, datetime_cols = local_column_groups(df)
    completeness = (1 - (summary.get("total_missing_values", 0) / max(summary.get("row_count", 1) * summary.get("column_count", 1), 1))) * 100
    duplicate_rate = (summary.get("duplicate_rows", 0) / max(summary.get("row_count", 1), 1)) * 100
    quality_score_value, quality_grade, quality_reasons = quality_score(summary)
    charts = _local_storyboard_chart_specs(df, dataset_id)
    kpis = _local_kpi_cards(df)[:6]
    anomalies = _local_anomaly_rows(df)
    anomalies.extend(_zscore_outlier_notes(df, numeric))
    anomalies = list(dict.fromkeys(anomalies))[:8]
    recommendations = _storyboard_recommendations(df)
    created_at = datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    dataset_name = _dataset_display_name(dataset_id)
    theme_snapshot = _storyboard_theme_snapshot(theme)

    summary_bullets = [
        f"{summary.get('row_count', 0):,} records across {summary.get('column_count', 0):,} columns were analyzed.",
        f"Data completeness is {completeness:.1f}% with duplicate rate at {duplicate_rate:.2f}%.",
        f"Detected {len(numeric):,} numeric fields, {len(categorical):,} categorical fields, and {len(datetime_cols):,} date-like fields.",
    ]
    risk_line = anomalies[0] if anomalies else "No major statistical anomalies detected in the selected dataset."
    summary_bullets.append(f"Key risk signal: {risk_line}")
    if datetime_cols and numeric:
        summary_bullets.append("Opportunity: use time-trend monitoring to detect directional shifts in core metrics.")
    elif categorical and numeric:
        summary_bullets.append("Opportunity: compare top segments to identify associated drivers of high and low performance.")
    summary_bullets.append("Recommended action: review quality and anomaly slides before publishing board decisions.")

    visual_slides = [
        {
            "slide_id": f"slide_data_visual_{dataset_id}",
            "title": "Data Visual Analysis",
            "section_type": "data_visual_analysis",
            "content": {
                "description": "Numeric distributions, category concentrations, and correlation/trend patterns from current dataset values.",
            },
            "charts": charts,
            "kpis": [],
            "insights": [
                "Visuals are auto-selected from available field types and rendered only when valid chart inputs exist.",
                "Correlation views indicate association strength and should not be interpreted as causal evidence.",
            ],
            "theme_snapshot": theme_snapshot,
        }
    ]

    business_insights: list[str] = []
    if categorical and numeric:
        top_col = categorical[0]
        top_counts = df[top_col].astype("string").fillna("Unknown").value_counts()
        if not top_counts.empty:
            business_insights.append(
                f"Top segment in {top_col} is '{top_counts.index[0]}' with {int(top_counts.iloc[0]):,} records, indicating concentration.")
    if datetime_cols and numeric:
        business_insights.append("A date-like field and numeric KPI are available, which supports trend-based monitoring.")
    if not business_insights:
        business_insights.append("Segment and trend insights are limited by current field mix; add richer categorical or date fields if available.")
    business_insights.append("What/Why/Action: identify the largest segment shifts, review likely associated factors, then prioritize one corrective action owner.")

    recommendations_lines = [
        f"{item['recommendation']} Reason: {item['reason']} Expected impact: {item['expected_impact']} Confidence: {item['confidence']}"
        for item in recommendations
    ]

    storyboard_items: list[dict] = [
        {
            "slide_id": f"slide_cover_{dataset_id}",
            "title": "Executive Board Cover",
            "section_type": "cover",
            "content": {
                "report_title": st.session_state.get("branding", {}).get("report_title", "Executive Decision Intelligence Report"),
                "dataset_name": dataset_name,
                "rows": summary.get("row_count", 0),
                "columns": summary.get("column_count", 0),
                "generated_at": created_at,
                "theme": theme_snapshot.get("theme"),
                "description": "Auto-built executive storyboard generated from your selected dataset.",
            },
            "charts": [],
            "kpis": [],
            "insights": ["Board-ready package generated automatically from local dataframe evidence."],
            "theme_snapshot": theme_snapshot,
        },
        {
            "slide_id": f"slide_exec_summary_{dataset_id}",
            "title": "Executive Summary",
            "section_type": "executive_summary",
            "content": {"description": "What happened, risks, opportunities, and recommended immediate action."},
            "charts": [],
            "kpis": [],
            "insights": summary_bullets[:6],
            "theme_snapshot": theme_snapshot,
        },
        {
            "slide_id": f"slide_data_quality_{dataset_id}",
            "title": "Data Quality Snapshot",
            "section_type": "data_quality",
            "content": {
                "description": "Completeness, missingness, duplicate risk, and overall quality grade.",
                "completeness_pct": round(completeness, 2),
                "missing_values": int(summary.get("total_missing_values", 0)),
                "duplicate_rate_pct": round(duplicate_rate, 3),
                "quality_score": quality_score_value,
                "quality_grade": quality_grade,
            },
            "charts": [],
            "kpis": [],
            "insights": quality_reasons,
            "theme_snapshot": theme_snapshot,
        },
        {
            "slide_id": f"slide_kpi_{dataset_id}",
            "title": "KPI Overview",
            "section_type": "kpi_overview",
            "content": {"description": "Core operational KPIs auto-selected from available numeric and profile fields."},
            "charts": [],
            "kpis": kpis,
            "insights": [
                "KPI cards summarize dataset health and magnitude indicators for executive scan speed.",
            ],
            "theme_snapshot": theme_snapshot,
        },
        *visual_slides,
        {
            "slide_id": f"slide_business_visual_{dataset_id}",
            "title": "Business Visual Analysis",
            "section_type": "business_visual_analysis",
            "content": {"description": "Top and bottom segments, trend signal, and risk-opportunity interpretation framework."},
            "charts": charts[:3],
            "kpis": [],
            "insights": business_insights,
            "theme_snapshot": theme_snapshot,
        },
        {
            "slide_id": f"slide_anomalies_{dataset_id}",
            "title": "Anomalies and Risks",
            "section_type": "anomalies_risks",
            "content": {"description": "IQR and z-score anomaly checks, missing-value risks, and dominance warnings."},
            "charts": [],
            "kpis": [],
            "insights": anomalies or ["No major statistical anomalies detected in the selected dataset."],
            "theme_snapshot": theme_snapshot,
        },
        {
            "slide_id": f"slide_recommendations_{dataset_id}",
            "title": "Recommendations",
            "section_type": "recommendations",
            "content": {"description": "Business recommendations with reason, expected impact, and confidence."},
            "charts": [],
            "kpis": [],
            "insights": recommendations_lines,
            "recommendations": recommendations,
            "theme_snapshot": theme_snapshot,
        },
        {
            "slide_id": f"slide_closing_{dataset_id}",
            "title": "Closing and Next Actions",
            "section_type": "closing",
            "content": {
                "description": "Prioritized next steps and optional data limitation note.",
                "limitations": "Insights indicate associations from observed data and should be validated before high-impact commitments.",
            },
            "charts": [],
            "kpis": [],
            "insights": [
                "Assign owners for quality remediation and anomaly monitoring.",
                "Publish a KPI watchlist with thresholds and review cadence.",
                "Use trend and segment visuals for next steering committee review.",
            ],
            "theme_snapshot": theme_snapshot,
        },
    ]
    return storyboard_items
def _storyboard_recommendations(df: pd.DataFrame) -> list[dict]:
    numeric, categorical, datetime_cols = local_column_groups(df)
    quality_summary = local_summary(df)
    _, _, quality_reasons = quality_score(quality_summary)
    recommendations: list[dict] = [
        {
            "recommendation": "Address the highest-missing columns before publishing final board metrics.",
            "reason": quality_reasons[-1] if quality_reasons else "Missing values can bias segmented KPI comparisons.",
            "expected_impact": "Higher confidence in KPI and segment-level interpretation.",
            "confidence": "High",
        },
        {
            "recommendation": "Monitor concentration in top categories and review whether distribution is expected.",
            "reason": "Category dominance indicates risk concentration and can mask smaller segment behavior.",
            "expected_impact": "Better prioritization of segment-specific actions.",
            "confidence": "Medium",
        },
    ]
    if datetime_cols and numeric:
        recommendations.append(
            {
                "recommendation": "Establish a weekly trend watch for the lead numeric KPI.",
                "reason": "Time-aware monitoring helps detect directional changes earlier.",
                "expected_impact": "Faster response to emerging performance drift.",
                "confidence": "Medium",
            }
        )
    if len(numeric) >= 2:
        recommendations.append(
            {
                "recommendation": "Track correlated numeric fields together in the same dashboard lane.",
                "reason": "Correlation patterns suggest variables move together and should be reviewed as a bundle.",
                "expected_impact": "Stronger root-cause triage during variance reviews.",
                "confidence": "Medium",
            }
        )
    recommendations.append(
        {
            "recommendation": "Define a recurring anomaly review cadence with owners and response thresholds.",
            "reason": "Outlier and duplication signals indicate where data or process controls may need tightening.",
            "expected_impact": "Reduced risk of acting on unstable or noisy observations.",
            "confidence": "Medium",
        }
    )
    return recommendations[:5]
def _auto_storyboard_entries(df: pd.DataFrame, dataset_id: str) -> list[dict]:
    entries: list[dict] = []
    for index, item in enumerate(_local_default_figures(df)[:8], start=1):
        entries.append(
            {
                "chart_id": f"auto_visual_{index}",
                "title": item["title"],
                "suggested_chart_type": item["kind"],
                "business_meaning": "Auto-built from dashboard visuals for the current local dataset.",
                "short_ai_insight": "Use this slide as a board-ready visual checkpoint.",
                "why_useful": "It gives the report package a real visual even when no manual selection has been made.",
                "spec": {"chart_type": "auto", "title": item["title"], "auto_index": index - 1, "dataset_id": dataset_id},
            }
        )
    return entries
def _build_storyboard_kpis(schema: dict) -> list[dict]:
    """Build simple KPI cards from visual builder schema metrics."""
    kpis = []
    row_count = 0
    semantic = schema.get("semantic_layer", [])
    measures = schema.get("measures", [])
    for field in semantic:
        if field.get("semantic_role") in {"revenue_currency_column"}:
            kpis.append({"label": f"Total {field['name']}", "value": "—", "icon": "chart"})
        elif field.get("semantic_role") in {"percentage_ratio_column"}:
            kpis.append({"label": field["name"].replace("_", " ").title(), "value": "—", "icon": "metric"})
    if not kpis:
        if measures:
            for m in measures[:2]:
                kpis.append({"label": m["name"].replace("_", " ").title(), "value": "—", "icon": "metric"})
        else:
            kpis.append({"label": "Records", "value": f"{row_count:,}" if row_count else "—", "icon": "table"})
    return kpis[:4]
def _local_storyboard_chart_specs(df: pd.DataFrame, dataset_id: str) -> list[dict]:
    figures = _local_default_figures(df)
    chart_specs: list[dict] = []
    for index, item in enumerate(figures, start=1):
        figure = item.get("figure")
        if not isinstance(figure, go.Figure):
            continue
        chart_specs.append(
            {
                "chart_id": f"story_auto_{dataset_id}_{index}",
                "title": item.get("title", f"Auto visual {index}"),
                "kind": item.get("kind", "Auto visual"),
                "plotly": figure.to_dict(),
            }
        )
    return chart_specs
def _storyboard_figures_from_items(items: list[dict]) -> list[dict]:
    figures: list[dict] = []
    for slide in items:
        for chart in slide.get("charts", []):
            plotly_spec = chart.get("plotly", {})
            if not isinstance(plotly_spec, dict):
                continue
            data = plotly_spec.get("data", [])
            layout = plotly_spec.get("layout", {})
            if not data:
                continue
            try:
                fig = go.Figure(data=data, layout=layout)
            except Exception:
                continue
            figures.append(
                {
                    "title": chart.get("title") or slide.get("title", "Storyboard Visual"),
                    "figure": fig,
                    "kind": chart.get("kind", "Storyboard"),
                }
            )
    return figures
def _local_storyboard_export_files(items: list[dict], dataset_id: str) -> dict[str, bytes]:
    figures = _storyboard_figures_from_items(items)
    files: dict[str, bytes] = {
        "json": json.dumps({"dataset_id": dataset_id, "storyboard_items": items}, indent=2, default=str).encode("utf-8"),
    }

    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt

        deck = Presentation()
        for slide in items:
            ppt_slide = deck.slides.add_slide(deck.slide_layouts[5])
            title = slide.get("title", "Storyboard Slide")
            ppt_slide.shapes.title.text = str(title)[:90]

            y = Inches(1.1)
            insights = slide.get("insights", [])
            if insights:
                textbox = ppt_slide.shapes.add_textbox(Inches(0.7), y, Inches(5.8), Inches(2.1))
                tf = textbox.text_frame
                tf.clear()
                tf.text = "Summary"
                tf.paragraphs[0].font.size = Pt(16)
                for item in insights[:4]:
                    p = tf.add_paragraph()
                    p.text = f"- {str(item)[:160]}"
                    p.level = 0
                    p.font.size = Pt(12)

            chart = next((c for c in slide.get("charts", []) if c.get("plotly", {}).get("data")), None)
            if chart:
                try:
                    fig = go.Figure(data=chart["plotly"].get("data", []), layout=chart["plotly"].get("layout", {}))
                    png = _figure_png_bytes(fig)
                    ppt_slide.shapes.add_picture(io.BytesIO(png), Inches(6.4), Inches(1.1), width=Inches(6.2), height=Inches(5.6))
                except Exception:
                    pass

            if not insights and not chart:
                textbox = ppt_slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.0), Inches(1.5))
                tf = textbox.text_frame
                tf.text = str(slide.get("content", {}).get("description", "Storyboard content"))[:240]
                tf.paragraphs[0].font.size = Pt(16)

        ppt_buffer = io.BytesIO()
        deck.save(ppt_buffer)
        files["pptx"] = ppt_buffer.getvalue()
    except Exception:
        files["pptx"] = b""

    try:
        from reportlab.lib.pagesizes import landscape, letter
        from reportlab.lib.utils import ImageReader
        from reportlab.pdfgen import canvas

        pdf_buffer = io.BytesIO()
        pdf = canvas.Canvas(pdf_buffer, pagesize=landscape(letter))
        width, height = landscape(letter)
        for slide in items:
            pdf.setFont("Helvetica-Bold", 18)
            pdf.drawString(36, height - 40, str(slide.get("title", "Storyboard Slide"))[:100])
            pdf.setFont("Helvetica", 11)
            y = height - 70
            for line in slide.get("insights", [])[:5]:
                pdf.drawString(42, y, f"- {str(line)[:140]}")
                y -= 16

            chart = next((c for c in slide.get("charts", []) if c.get("plotly", {}).get("data")), None)
            if chart:
                try:
                    fig = go.Figure(data=chart["plotly"].get("data", []), layout=chart["plotly"].get("layout", {}))
                    png = _figure_png_bytes(fig)
                    pdf.drawImage(ImageReader(io.BytesIO(png)), 42, 46, width=width - 84, height=height - 180, preserveAspectRatio=True, anchor="c")
                except Exception:
                    pass
            pdf.showPage()
        pdf.save()
        files["pdf"] = pdf_buffer.getvalue()
    except Exception:
        files["pdf"] = b""

    return files
